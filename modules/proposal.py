"""
Customer Proposal Generator — 38 Degrees North branded PPTX.

Design system derived from:
  - 38DN official theme (Color Scheme "38N", Font Scheme "Century Gothic")
  - McKinsey presentation conventions (action titles, Pyramid Principle, source lines)
  - BCG slide structure (bumper/takeaway boxes, exhibit numbering)
  - JP Morgan pitchbook formatting (clean tables, data-first, institutional tone)

Every data slide follows a strict four-zone layout:
  ┌─ Navy accent rule (4 pt) ─────────────────────────────┐
  │  ACTION TITLE — a complete sentence stating the insight│
  │  Subtitle — supporting context in gray                 │
  ├────────────────────────────────────────────────────────┤
  │  TAKEAWAY BOX — light grey (#F0F0F2) with green bar   │
  ├────────────────────────────────────────────────────────┤
  │                                                        │
  │                  BODY CONTENT                           │
  │         (tables / KPIs / charts / text)                │
  │                                                        │
  ├────────────────────────────────────────────────────────┤
  │  Source: ...                   38DN | CONFIDENTIAL  p# │
  └────────────────────────────────────────────────────────┘
"""

from __future__ import annotations

import os
from io import BytesIO

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, ChartData

from .billing import BillingResult

# ═══════════════════════════════════════════════════════════════════════════
# 38DN BRAND TOKENS  (from theme "38N" extracted from MBR Template)
# ═══════════════════════════════════════════════════════════════════════════
FONT       = "Century Gothic"          # majorFont + minorFont

DK1        = RGBColor(0x05, 0x0D, 0x25)  # primary navy
DK2        = RGBColor(0x21, 0x2B, 0x48)  # secondary dark
LT_GRAY    = RGBColor(0xF0, 0xF0, 0xF2)  # light grey (callout bg)
ACCENT1    = RGBColor(0x45, 0xA7, 0x50)  # 38DN green
ACCENT3    = RGBColor(0x51, 0x84, 0x84)  # teal
ACCENT4    = RGBColor(0x1D, 0x6F, 0xA9)  # blue

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
GRAY50     = RGBColor(0x66, 0x66, 0x66)
GRAY70     = RGBColor(0x99, 0x99, 0x99)
GRAY85     = RGBColor(0xCC, 0xCC, 0xCC)

TBL_HDR    = DK2                          # table header bg
TBL_BORDER = RGBColor(0xDE, 0xE2, 0xE6)  # table cell border
TBL_ALT    = RGBColor(0xF7, 0xF8, 0xFA)  # alternate row bg

# Slide dimensions — standard 16:9 (matches Board Deck)
SW = Inches(13.333)
SH = Inches(7.5)

# Layout grid  (1" margins like McKinsey spec)
ML = Inches(0.75)          # margin left
MR = Inches(0.75)
CW = Inches(11.833)        # content width (SW - ML - MR)

# Vertical zones
RULE_Y   = Emu(0)          # top accent rule
RULE_H   = Pt(4)           # thin accent line
TITLE_Y  = Inches(0.22)    # action title top
TITLE_H  = Inches(0.50)
SUB_Y    = Inches(0.68)    # subtitle
SUB_H    = Inches(0.30)
BODY_Y   = Inches(1.05)    # start of body content
SRC_Y    = Inches(6.85)    # source line
FTR_Y    = Inches(7.10)    # footer

# Logo paths
LOGO_MARK  = "logo.png"      # 38 circle mark (white bg)
LOGO_WORD  = "38NORTH.png"   # horizontal wordmark (transparent)

def _logo_path(filename):
    candidates = [
        os.path.join(os.path.dirname(__file__), "..", "assets", filename),
        os.path.join(os.path.expanduser("~"),
                     "38 Degrees North", "38DN - Documents", "General",
                     "38DN Logos", filename),
    ]
    for p in candidates:
        if os.path.isfile(p):
            return p
    return None


# ═══════════════════════════════════════════════════════════════════════════
# SHAPE PRIMITIVES
# ═══════════════════════════════════════════════════════════════════════════

def _rect(sl, l, t, w, h, fill=None, line=None, line_w=None):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.shadow.inherit = False
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = line_w or Pt(0.5)
    return s

def _txt(sl, l, t, w, h, text="", sz=Pt(12), bold=False, italic=False,
         color=DK1, align=PP_ALIGN.LEFT, font=FONT, wrap=True):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = sz
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return tb

def _multi(sl, l, t, w, h, runs, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for rd in runs:
        r = p.add_run()
        r.text = rd.get("t", "")
        r.font.name = rd.get("f", FONT)
        r.font.size = rd.get("sz", Pt(12))
        r.font.bold = rd.get("b", False)
        r.font.italic = rd.get("i", False)
        r.font.color.rgb = rd.get("c", DK1)
    return tb

def _bullets(sl, l, t, w, items, sz=Pt(10), color=DK1, spacing=Pt(6)):
    tb = sl.shapes.add_textbox(l, t, w, Inches(len(items) * 0.32))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = spacing
        r = p.add_run(); r.text = f"\u2022  {item}"
        r.font.name = FONT; r.font.size = sz; r.font.color.rgb = color


# ═══════════════════════════════════════════════════════════════════════════
# STRUCTURAL ELEMENTS  (McKinsey / BCG / JPM conventions)
# ═══════════════════════════════════════════════════════════════════════════

def _accent_rule(sl):
    """4-pt navy rule across the very top — thin, elegant, institutional."""
    _rect(sl, Emu(0), RULE_Y, SW, Pt(4), fill=DK1)

def _action_title(sl, text, exhibit=None):
    """Action title: a sentence that IS the conclusion (McKinsey convention).
    Optional exhibit number prefix (BCG convention)."""
    prefix = f"Exhibit {exhibit}:  " if exhibit else ""
    _txt(sl, ML, TITLE_Y, CW, TITLE_H, text=prefix + text,
         sz=Pt(20), bold=True, color=DK1)

def _subtitle(sl, text):
    """Gray context line under the action title."""
    _txt(sl, ML, SUB_Y, CW, SUB_H, text=text,
         sz=Pt(11), color=GRAY50)

def _takeaway(sl, text, y=None):
    """BCG-style 'so what' bumper box — light grey with green accent bar."""
    ty = y or BODY_Y
    _rect(sl, ML, ty, CW, Inches(0.48), fill=LT_GRAY)
    _rect(sl, ML, ty, Pt(4), Inches(0.48), fill=ACCENT1)
    _txt(sl, ML + Inches(0.18), ty + Inches(0.07), CW - Inches(0.3), Inches(0.34),
         text=text, sz=Pt(10.5), bold=True, color=DK1)

def _source(sl, text, y=None):
    """Source/footnote line — lower left, small italic (McKinsey convention)."""
    _txt(sl, ML, y or SRC_Y, CW, Inches(0.22), text=f"Source: {text}",
         sz=Pt(7.5), italic=True, color=GRAY70)

def _footer(sl, pg, total):
    """Footer: company | CONFIDENTIAL + page number (right-aligned)."""
    _multi(sl, ML, FTR_Y, Inches(4), Inches(0.25),
           [{"t": "38 DEGREES NORTH", "sz": Pt(7), "b": True, "c": GRAY70},
            {"t": "  |  CONFIDENTIAL", "sz": Pt(7), "c": GRAY70}])
    _txt(sl, Inches(11.5), FTR_Y, Inches(1.2), Inches(0.25),
         text=f"{pg} / {total}", sz=Pt(7), color=GRAY70, align=PP_ALIGN.RIGHT)

def _divider_line(sl, y, width=None, x=None):
    """Thin horizontal rule for section separation."""
    _rect(sl, x if x is not None else ML, y, width or CW, Pt(1), fill=GRAY85)


# ═══════════════════════════════════════════════════════════════════════════
# KPI TILES  (institutional metric card — no background, just accent + data)
# ═══════════════════════════════════════════════════════════════════════════

def _kpi(sl, l, t, w, value, label, accent=DK1, val_color=DK1):
    _rect(sl, l, t, w, Pt(3), fill=accent)
    _txt(sl, l, t + Inches(0.12), w, Inches(0.40), text=value,
         sz=Pt(24), bold=True, color=val_color, align=PP_ALIGN.CENTER)
    _txt(sl, l, t + Inches(0.52), w, Inches(0.22), text=label,
         sz=Pt(9), color=GRAY50, align=PP_ALIGN.CENTER)

def _kpi_row(sl, tiles, y, tile_w=Inches(2.7)):
    n = len(tiles)
    gap = (CW - n * tile_w) / max(n - 1, 1)
    for i, t in enumerate(tiles):
        x = ML + i * (tile_w + gap)
        _kpi(sl, x, y, tile_w, **t)


# ═══════════════════════════════════════════════════════════════════════════
# TABLE  (JP Morgan pitchbook style — clean, tight, institutional)
# ═══════════════════════════════════════════════════════════════════════════

def _table(sl, l, t, w, col_ws, hdrs, rows, bold_last=False, sz=Pt(9)):
    nr = len(rows) + 1; nc = len(hdrs)
    row_h = 0.28 if sz >= Pt(9) else 0.19
    sh = sl.shapes.add_table(nr, nc, l, t, w, Inches(row_h * nr))
    tbl = sh.table
    for i, cw in enumerate(col_ws):
        tbl.columns[i].width = cw

    # Header
    for j, h in enumerate(hdrs):
        c = tbl.cell(0, j); c.text = h
        _cfmt(c, bold=True, sz=sz, color=WHITE, bg=TBL_HDR,
              al=PP_ALIGN.CENTER)

    # Data
    for i, row in enumerate(rows):
        bg = TBL_ALT if i % 2 == 1 else None
        is_last = bold_last and i == len(rows) - 1
        for j, v in enumerate(row):
            c = tbl.cell(i+1, j); c.text = str(v)
            _cfmt(c, bold=is_last, sz=sz, color=DK1, bg=bg,
                  al=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT)

    # Thin borders via XML
    _apply_table_borders(tbl, nr, nc)
    return sh

def _cfmt(cell, bold=False, sz=Pt(9), color=DK1, bg=None, al=PP_ALIGN.RIGHT):
    for p in cell.text_frame.paragraphs:
        p.alignment = al
        for r in p.runs:
            r.font.name = FONT; r.font.size = sz
            r.font.bold = bold; r.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)

def _apply_table_borders(tbl, nr, nc):
    """Apply thin #DEE2E6 borders to all cells via XML manipulation."""
    from pptx.oxml.ns import qn
    from lxml import etree
    border_color = "DEE2E6"
    border_w = "6350"  # 0.5pt in EMU
    for i in range(nr):
        for j in range(nc):
            cell = tbl.cell(i, j)
            tc = cell._tc
            tcPr = tc.find(qn('a:tcPr'))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn('a:tcPr'))
            for side in ['lnL', 'lnR', 'lnT', 'lnB']:
                ln = etree.SubElement(tcPr, qn(f'a:{side}'))
                ln.set('w', border_w)
                ln.set('cap', 'flat')
                ln.set('cmpd', 'sng')
                sf = etree.SubElement(ln, qn('a:solidFill'))
                sc = etree.SubElement(sf, qn('a:srgbClr'))
                sc.set('val', border_color)


# ═══════════════════════════════════════════════════════════════════════════
# STEP LIST  (numbered process — How It Works / Next Steps)
# ═══════════════════════════════════════════════════════════════════════════

def _step(sl, l, t, num, title, desc, circ_c=ACCENT1, dark=False):
    tc = WHITE if dark else DK1; dc = GRAY70 if dark else GRAY50
    c = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(0.38), Inches(0.38))
    c.fill.solid(); c.fill.fore_color.rgb = circ_c; c.line.fill.background()
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num); r.font.name = FONT
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
    _txt(sl, l + Inches(0.52), t - Inches(0.01), Inches(4.8), Inches(0.26),
         text=title, sz=Pt(11.5), bold=True, color=tc)
    _txt(sl, l + Inches(0.52), t + Inches(0.24), Inches(4.8), Inches(0.28),
         text=desc, sz=Pt(9), color=dc)


# ═══════════════════════════════════════════════════════════════════════════
# FORMAT HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def _fd(v):
    if abs(v) >= 1e6: return f"${v/1e6:,.2f}M"
    if abs(v) >= 1e3: return f"${v/1e3:,.0f}K"
    return f"${v:,.0f}"

def _fk(v):
    if abs(v) >= 1e6: return f"{v/1e6:,.1f} GWh"
    if abs(v) >= 1e3: return f"{v/1e3:,.0f} MWh"
    return f"{v:,.0f} kWh"


# ═══════════════════════════════════════════════════════════════════════════
# NATIVE PPTX CHARTS  (Excel-linked, editable in PowerPoint)
# ═══════════════════════════════════════════════════════════════════════════

def _style_chart_font(chart, font_name=FONT, sz_pt=800):
    """Apply Century Gothic at *sz_pt* (hundredths of a point) to all chart
    text: legend, axis tick labels, and axis titles.  Default 800 = 8 pt."""
    from pptx.oxml.ns import qn
    from lxml import etree

    def _apply(root_elem):
        if root_elem is None:
            return
        for tag in ('a:rPr', 'a:defRPr'):
            for node in root_elem.findall('.//' + qn(tag)):
                node.set('sz', str(sz_pt))
                latin = node.find(qn('a:latin'))
                if latin is None:
                    latin = etree.SubElement(node, qn('a:latin'))
                latin.set('typeface', font_name)

    def _ensure_defRPr(elem, tag_local='c:txPr'):
        """If the element has no text properties, inject a minimal txPr with
        a defRPr so _apply can stamp the font on it."""
        txPr = elem.find(qn(tag_local))
        if txPr is None:
            txPr = etree.SubElement(elem, qn(tag_local))
        bodyPr = txPr.find(qn('a:bodyPr'))
        if bodyPr is None:
            etree.SubElement(txPr, qn('a:bodyPr'))
        p = txPr.find(qn('a:p'))
        if p is None:
            p = etree.SubElement(txPr, qn('a:p'))
        pPr = p.find(qn('a:pPr'))
        if pPr is None:
            pPr = etree.SubElement(p, qn('a:pPr'))
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is None:
            etree.SubElement(pPr, qn('a:defRPr'))

    # Legend
    if chart.has_legend:
        _ensure_defRPr(chart.legend._element)
        _apply(chart.legend._element)

    # Axes — tick labels + titles
    for axis in (chart.value_axis, chart.category_axis):
        _ensure_defRPr(axis._element)
        _apply(axis._element)
        if axis.has_title:
            _apply(axis.axis_title._element)


def _add_prod_load_chart(sl, left, top, width, height, monthly_summary):
    """Add a native grouped column chart: Load vs Solar by month."""
    months = ["Jan","Feb","Mar","Apr","May","Jun",
              "Jul","Aug","Sep","Oct","Nov","Dec"]
    ms = monthly_summary
    n = min(12, len(ms))

    chart_data = CategoryChartData()
    chart_data.categories = months[:n]
    chart_data.add_series("Site Load (MWh)",
        [ms.iloc[i].get("load_kwh", 0) / 1e3 for i in range(n)])
    chart_data.add_series("Solar Production (MWh)",
        [ms.iloc[i].get("solar_kwh", 0) / 1e3 for i in range(n)])

    chart_frame = sl.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data)
    chart = chart_frame.chart
    chart.chart_style = 2                  # minimal built-in style
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Style series colours
    plot = chart.plots[0]
    plot.gap_width = 80
    s0 = plot.series[0]; s0.format.fill.solid(); s0.format.fill.fore_color.rgb = DK2
    s1 = plot.series[1]; s1.format.fill.solid(); s1.format.fill.fore_color.rgb = ACCENT1

    # Y-axis
    va = chart.value_axis
    va.has_title = True; va.axis_title.text_frame.paragraphs[0].text = "MWh"
    va.major_gridlines.format.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    va.major_gridlines.format.line.width = Pt(0.4)
    va.format.line.color.rgb = GRAY85; va.format.line.width = Pt(0.5)

    # Category axis
    ca = chart.category_axis
    ca.format.line.color.rgb = GRAY85; ca.format.line.width = Pt(0.5)

    _style_chart_font(chart)
    return chart_frame


def _add_hedge_chart(sl, left, top, width, height,
                     baseline_bill, ppa_bill_yr1, ppa_esc_pct, term,
                     solar_kwh=0, ppa_rate_val=None):
    """Add a native line chart: PPA cost vs utility at 4%, 7%, 10%."""
    years_list = list(range(1, term + 1))
    years = np.array(years_list)

    # Utility scenarios
    util_4 = baseline_bill * (1.04 ** (years - 1))
    util_7 = baseline_bill * (1.07 ** (years - 1))
    util_10 = baseline_bill * (1.10 ** (years - 1))

    # PPA total cost
    grid_residual_yr1 = ppa_bill_yr1
    if ppa_rate_val and solar_kwh:
        ppa_energy_yr1 = ppa_rate_val * solar_kwh
        grid_residual_yr1 = ppa_bill_yr1 - ppa_energy_yr1
        ppa_energy = ppa_energy_yr1 * ((1 + ppa_esc_pct / 100) ** (years - 1))
    else:
        ppa_energy = np.zeros_like(years, dtype=float)
    grid_residual = grid_residual_yr1 * (1.04 ** (years - 1))
    ppa_total = grid_residual + ppa_energy

    chart_data = CategoryChartData()
    chart_data.categories = [str(y) for y in years_list]
    chart_data.add_series("Utility Only @ 10%/yr", (util_10 / 1e3).tolist())
    chart_data.add_series("Utility Only @ 7%/yr", (util_7 / 1e3).tolist())
    chart_data.add_series("Utility Only @ 4%/yr", (util_4 / 1e3).tolist())
    chart_data.add_series("With PPA + Solar", (ppa_total / 1e3).tolist())

    chart_frame = sl.shapes.add_chart(
        XL_CHART_TYPE.LINE, left, top, width, height, chart_data)
    chart = chart_frame.chart
    chart.chart_style = 2
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.smooth = False

    # Style series: 10% blue dotted, 7% navy solid, 4% teal dashed, PPA green thick
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    colors = [ACCENT4, DK1, ACCENT3, ACCENT1]
    widths = [Pt(1.3), Pt(1.8), Pt(1.3), Pt(2.2)]
    dashes = [MSO_LINE_DASH_STYLE.DASH_DOT, MSO_LINE_DASH_STYLE.SOLID,
              MSO_LINE_DASH_STYLE.DASH, MSO_LINE_DASH_STYLE.SOLID]
    for i, s in enumerate(plot.series):
        s.format.line.color.rgb = colors[i]
        s.format.line.width = widths[i]
        s.format.line.dash_style = dashes[i]
        s.smooth = False

    # Y-axis
    va = chart.value_axis
    va.has_title = True; va.axis_title.text_frame.paragraphs[0].text = "Annual Cost ($K)"
    va.major_gridlines.format.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    va.major_gridlines.format.line.width = Pt(0.4)
    va.format.line.color.rgb = GRAY85; va.format.line.width = Pt(0.5)

    # Category axis
    ca = chart.category_axis
    ca.has_title = True; ca.axis_title.text_frame.paragraphs[0].text = "Contract Year"
    ca.format.line.color.rgb = GRAY85; ca.format.line.width = Pt(0.5)
    from pptx.enum.chart import XL_TICK_LABEL_POSITION
    ca.tick_label_position = XL_TICK_LABEL_POSITION.LOW

    _style_chart_font(chart)
    return chart_frame


def _add_savings_components_chart(sl, left, top, width, height,
                                  proj_df, yr1_base_energy, yr1_base_demand,
                                  yr1_base_fixed, rate_esc_pct):
    """Add a native stacked column chart: savings by component per year."""
    df = proj_df.copy()
    if "Year" not in df.columns:
        df.insert(0, "Year", range(1, len(df) + 1))
    years = df["Year"].values

    rate_factors = 1.0 + (rate_esc_pct / 100) * (years - 1)
    base_energy = yr1_base_energy * rate_factors
    base_demand = np.full_like(years, yr1_base_demand, dtype=float)

    solar_energy = df["Energy ($)"].values.astype(float)
    solar_demand = df["Demand ($)"].values.astype(float)
    export_cr = df["Export Credit ($)"].values.astype(float)
    nbc = df["NBC ($)"].values.astype(float) if "NBC ($)" in df.columns else np.zeros(len(df))

    energy_sav = np.maximum(base_energy - solar_energy, 0)
    demand_sav = np.maximum(base_demand - solar_demand, 0)
    export_sav = np.abs(export_cr)

    chart_data = CategoryChartData()
    chart_data.categories = [str(int(y)) for y in years]
    chart_data.add_series("Energy Savings", (energy_sav / 1e3).tolist())
    chart_data.add_series("Demand Savings", (demand_sav / 1e3).tolist())
    chart_data.add_series("Export Credits", (export_sav / 1e3).tolist())
    if np.any(nbc > 0):
        chart_data.add_series("NBC Charges", (-nbc / 1e3).tolist())

    chart_frame = sl.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, chart_data)
    chart = chart_frame.chart
    chart.chart_style = 2
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.gap_width = 60
    series_colors = [ACCENT1, ACCENT3, ACCENT4, RGBColor(0xCC, 0x44, 0x44)]
    for i, s in enumerate(plot.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = series_colors[min(i, len(series_colors)-1)]

    va = chart.value_axis
    va.has_title = True; va.axis_title.text_frame.paragraphs[0].text = "Annual Savings ($K)"
    va.major_gridlines.format.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    va.major_gridlines.format.line.width = Pt(0.4)
    va.format.line.color.rgb = GRAY85; va.format.line.width = Pt(0.5)

    ca = chart.category_axis
    ca.has_title = True; ca.axis_title.text_frame.paragraphs[0].text = "Contract Year"
    ca.format.line.color.rgb = GRAY85; ca.format.line.width = Pt(0.5)

    _style_chart_font(chart)
    return chart_frame


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def _slide_cover(prs, name, addr, acct, date_str, total):
    """Title slide — white bg, logomark top-left, clean institutional."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # Thin navy rule at top
    _rect(sl, Emu(0), Emu(0), SW, Pt(5), fill=DK1)

    # Logo mark (top left)
    lp = _logo_path(LOGO_MARK)
    if lp:
        sl.shapes.add_picture(lp, Inches(0.6), Inches(0.35), Inches(0.65))

    # Thin navy line below title area
    _rect(sl, ML, Inches(3.95), CW, Pt(1.5), fill=DK1)

    # Title
    _txt(sl, ML, Inches(2.4), CW, Inches(0.7),
         text="Behind-the-Meter Solar + Storage",
         sz=Pt(30), bold=True, color=DK1, align=PP_ALIGN.LEFT)
    _txt(sl, ML, Inches(3.1), CW, Inches(0.5),
         text="Investment Summary",
         sz=Pt(22), color=ACCENT1, align=PP_ALIGN.LEFT)

    # Customer details below the line
    _txt(sl, ML, Inches(4.25), CW, Inches(0.40),
         text=name, sz=Pt(16), bold=True, color=DK1)

    sub = addr
    if acct:
        sub += f"  |  {acct}"
    if sub:
        _txt(sl, ML, Inches(4.70), CW, Inches(0.30),
             text=sub, sz=Pt(10), color=GRAY50)

    # Date + Confidential at bottom
    _txt(sl, ML, Inches(6.6), CW, Inches(0.25),
         text=f"Confidential  |  {date_str}",
         sz=Pt(9), color=GRAY70)

    # Wordmark bottom-right
    wp = _logo_path(LOGO_WORD)
    if wp:
        sl.shapes.add_picture(wp, Inches(10.0), Inches(6.45), Inches(2.5))


def _slide_exec_summary(prs, pg, total, name, result, tariff, utility,
                        sys_kw, batt_kwh, ppa_rate, esc_pct, term,
                        customer_savings_pct=None, customer_savings_pct_2=None,
                        nem_regime_1=None, nem_regime_2=None):
    """Executive Summary — two-column financial + technical overview."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _accent_rule(sl)

    bill = result.annual_bill_without_solar
    sav = result.annual_savings
    pct = customer_savings_pct if customer_savings_pct is not None else result.savings_pct
    offset = result.annual_solar_kwh / result.annual_load_kwh * 100 if result.annual_load_kwh else 0

    _action_title(sl, "Executive Summary")
    _subtitle(sl, f"{name}  |  {utility} {tariff}  |  {sys_kw:,.0f} kW-DC Solar"
              + (f" + {batt_kwh:,.0f} kWh BESS" if batt_kwh > 0 else ""))
    _takeaway(sl,
        f"Day-1 savings of {_fd(sav)} ({pct:.1f}%) with zero upfront capital. "
        f"System offsets {offset:.0f}% of annual load and provides long-term rate certainty.")

    # ── Left column: Financial ──
    cx1 = ML; cx2 = ML + Inches(6.2)
    sy = Inches(1.80)

    _txt(sl, cx1, sy, Inches(5.5), Inches(0.22),
         text="FINANCIAL OVERVIEW", sz=Pt(11), bold=True, color=DK1)
    _divider_line(sl, sy + Inches(0.27), Inches(5.5))

    peak = result.monthly_summary["peak_demand_kw"].max() if "peak_demand_kw" in result.monthly_summary.columns else 0
    fin = [
        ("Current Annual Utility Bill", _fd(bill)),
        ("Projected Year 1 Cost (w/ Solar)", _fd(result.annual_bill_with_solar)),
        ("Year 1 Net Savings", _fd(sav)),
    ]
    if customer_savings_pct_2 is not None and nem_regime_1 and nem_regime_2:
        fin.append((f"{nem_regime_1} Savings Target", f"{pct:.1f}%"))
        fin.append((f"{nem_regime_2} Savings Target", f"{customer_savings_pct_2:.1f}%"))
    else:
        fin.append(("Savings Rate", f"{pct:.1f}%"))
    if ppa_rate is not None:
        fin.append(("PPA Rate (Year 1)", f"${ppa_rate:.3f}/kWh"))
    if esc_pct is not None:
        fin.append(("PPA Escalator", f"{esc_pct:.1f}%/yr"))
    fin.append(("Contract Term", f"{term} years"))
    fin.append(("Upfront Cost", "$0"))

    for i, (lbl, val) in enumerate(fin):
        y = sy + Inches(0.40 + i * 0.36)
        _txt(sl, cx1 + Inches(0.1), y, Inches(3.5), Inches(0.26),
             text=lbl, sz=Pt(10), color=GRAY50)
        _txt(sl, cx1 + Inches(3.7), y, Inches(1.7), Inches(0.26),
             text=val, sz=Pt(10), bold=True, color=DK1, align=PP_ALIGN.RIGHT)

    # ── Right column: Technical ──
    _txt(sl, cx2, sy, Inches(5.5), Inches(0.22),
         text="SYSTEM & ENERGY", sz=Pt(11), bold=True, color=DK1)
    _divider_line(sl, sy + Inches(0.27), Inches(5.5))

    yld = result.annual_solar_kwh / sys_kw if sys_kw else 0
    sc = ((result.annual_solar_kwh - result.annual_export_kwh) /
          result.annual_solar_kwh * 100 if result.annual_solar_kwh else 0)

    tech = [
        ("System Size (DC)", f"{sys_kw:,.0f} kW"),
        ("Annual Production", _fk(result.annual_solar_kwh)),
        ("Annual Load", _fk(result.annual_load_kwh)),
        ("Solar Offset", f"{offset:.0f}%"),
        ("Specific Yield", f"{yld:,.0f} kWh/kW"),
        ("Self-Consumption", f"{sc:.0f}%"),
        ("Peak Demand", f"{peak:,.0f} kW"),
    ]
    if batt_kwh > 0:
        tech.append(("Battery Storage", f"{batt_kwh:,.0f} kWh"))

    for i, (lbl, val) in enumerate(tech):
        y = sy + Inches(0.40 + i * 0.36)
        _txt(sl, cx2 + Inches(0.1), y, Inches(3.5), Inches(0.26),
             text=lbl, sz=Pt(10), color=GRAY50)
        _txt(sl, cx2 + Inches(3.7), y, Inches(1.7), Inches(0.26),
             text=val, sz=Pt(10), bold=True, color=DK1, align=PP_ALIGN.RIGHT)

    _source(sl, "38DN pv-rate-sim billing model; 8,760-hour simulation")
    _footer(sl, pg, total)


def _slide_current_cost(prs, pg, total, ex, name, result, tariff, utility):
    """Baseline electricity cost analysis."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _accent_rule(sl)

    bill = result.annual_bill_without_solar
    peak = result.monthly_summary["peak_demand_kw"].max() if "peak_demand_kw" in result.monthly_summary.columns else 0
    blended = bill / result.annual_load_kwh if result.annual_load_kwh else 0

    _action_title(sl,
        f"{name} spends {_fd(bill)}/yr on electricity ({tariff})", exhibit=ex)
    _subtitle(sl,
        f"{utility}  |  {result.annual_load_kwh/1e3:,.0f} MWh annual load  |  "
        f"{peak:,.0f} kW peak demand  |  ${blended:.4f}/kWh blended rate")
    _takeaway(sl,
        f"All-in blended rate of ${blended:.4f}/kWh. Demand charges represent a "
        f"significant cost component. Rising utility rates compound exposure annually.")

    # KPI row
    tiles = [
        {"value": _fd(bill), "label": "Annual Utility Bill", "accent": DK1},
        {"value": f"{peak:,.0f} kW", "label": "Peak Demand", "accent": ACCENT3},
        {"value": f"${blended:.4f}", "label": "Blended $/kWh", "accent": ACCENT1},
        {"value": _fk(result.annual_load_kwh), "label": "Annual Load", "accent": ACCENT4},
    ]
    _kpi_row(sl, tiles, Inches(1.80))

    # Monthly table
    _txt(sl, ML, Inches(2.75), CW, Inches(0.20),
         text="MONTHLY COST BREAKDOWN  (WITHOUT SOLAR)", sz=Pt(11), bold=True, color=DK1)
    _divider_line(sl, Inches(2.98))

    months = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
    ms = result.monthly_summary
    hdrs = ["Month", "Load (MWh)", "Peak kW", "Energy ($)", "Demand ($)", "Fixed ($)", "Total ($)"]

    rows_data = []
    use_baseline = result.monthly_baseline_details is not None
    for i in range(min(12, len(ms))):
        r = ms.iloc[i]
        if use_baseline:
            d = result.monthly_baseline_details[i]
            rows_data.append([
                months[i],
                f"{r.get('load_kwh',0)/1e3:,.1f}",
                f"{r.get('peak_demand_kw',0):,.0f}",
                f"${d.get('energy',0):,.0f}",
                f"${d.get('demand',0):,.0f}",
                f"${d.get('fixed',0):,.0f}",
                f"${d.get('total',0):,.0f}",
            ])
        else:
            rows_data.append([
                months[i],
                f"{r.get('load_kwh',0)/1e3:,.1f}",
                f"{r.get('peak_demand_kw',0):,.0f}",
                f"${r.get('energy_cost',0):,.0f}",
                f"${r.get('total_demand_charge',0):,.0f}",
                f"${r.get('fixed_charge',0):,.0f}",
                f"${r.get('net_bill',0):,.0f}",
            ])

    cws = [Inches(0.85), Inches(1.25), Inches(1.0), Inches(1.4), Inches(1.4), Inches(1.2), Inches(1.4)]
    _table(sl, ML, Inches(3.05), Inches(8.5), cws, hdrs, rows_data)

    _source(sl, f"38DN billing simulation; {utility} {tariff} rate schedule; 8,760-hour model")
    _footer(sl, pg, total)


def _slide_system(prs, pg, total, ex, sys_kw, dc_ac, batt_kwh, batt_kw,
                  ppa, esc, term, tariff, new_tariff):
    """Proposed system configuration and deal terms."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _accent_rule(sl)

    ac_kw = sys_kw / dc_ac if dc_ac > 0 else sys_kw
    t = f"Proposed {sys_kw:,.0f} kW solar"
    if batt_kwh > 0: t += f" + {batt_kwh:,.0f} kWh storage"
    t += " at zero upfront cost"

    _action_title(sl, t, exhibit=ex)
    _subtitle(sl, "System sized to optimize load offset and demand charge reduction under current tariff structure")

    # PV panel
    pw = Inches(5.2); ph = Inches(2.8); py = Inches(1.25)
    _rect(sl, ML, py, pw, ph, fill=None, line=TBL_BORDER, line_w=Pt(1))
    _rect(sl, ML, py, pw, Pt(3.5), fill=ACCENT1)
    _txt(sl, ML + Inches(0.25), py + Inches(0.15), Inches(4), Inches(0.25),
         text="SOLAR ARRAY (PV)", sz=Pt(9), bold=True, color=DK1)

    specs_pv = [("DC Capacity", f"{sys_kw:,.0f} kW"), ("AC Interconnection", f"{ac_kw:,.0f} kW"),
                ("DC:AC Ratio", f"{dc_ac:.1f}:1")]
    for i, (lbl, val) in enumerate(specs_pv):
        y = py + Inches(0.55 + i * 0.48)
        _txt(sl, ML + Inches(0.3), y, Inches(2.2), Inches(0.25), text=lbl,
             sz=Pt(10.5), color=GRAY50)
        _txt(sl, ML + Inches(2.6), y, Inches(2.2), Inches(0.25), text=val,
             sz=Pt(10.5), bold=True, color=DK1)

    # Battery panel
    if batt_kwh > 0:
        bx = ML + Inches(5.7)
        _rect(sl, bx, py, pw, ph, fill=None, line=TBL_BORDER, line_w=Pt(1))
        _rect(sl, bx, py, pw, Pt(3.5), fill=ACCENT3)
        _txt(sl, bx + Inches(0.25), py + Inches(0.15), Inches(4), Inches(0.25),
             text="BATTERY STORAGE (BESS)", sz=Pt(9), bold=True, color=DK1)
        dur = batt_kwh / batt_kw if batt_kw > 0 else 0
        specs_b = [("Energy Capacity", f"{batt_kwh:,.0f} kWh"),
                   ("Power Rating", f"{batt_kw:,.0f} kW"),
                   ("Duration", f"{dur:.1f} hours")]
        for i, (lbl, val) in enumerate(specs_b):
            y = py + Inches(0.55 + i * 0.48)
            _txt(sl, bx + Inches(0.3), y, Inches(2.2), Inches(0.25), text=lbl,
                 sz=Pt(10.5), color=GRAY50)
            _txt(sl, bx + Inches(2.6), y, Inches(2.2), Inches(0.25), text=val,
                 sz=Pt(10.5), bold=True, color=DK1)

    # Deal terms bar
    by = Inches(4.4)
    _rect(sl, ML, by, CW, Inches(0.48), fill=DK1)
    parts = []
    if ppa is not None:
        parts.append({"t": f"PPA Rate: ${ppa:.3f}/kWh", "c": ACCENT1, "b": True})
    if esc is not None:
        parts.append({"t": f"Escalator: {esc:.1f}%/yr", "c": WHITE, "b": True})
    parts += [{"t": f"Term: {term} years", "c": WHITE, "b": True},
              {"t": "Upfront Cost: $0", "c": WHITE, "b": True},
              {"t": f"Tariff: {new_tariff or tariff}", "c": WHITE, "b": True}]
    runs = []
    for i, p in enumerate(parts):
        if i > 0:
            runs.append({"t": "   |   ", "sz": Pt(9.5), "c": GRAY70})
        runs.append({"t": p["t"], "sz": Pt(9.5), "b": p["b"], "c": p["c"]})
    _multi(sl, ML + Inches(0.15), by + Inches(0.07), CW - Inches(0.3), Inches(0.34),
           runs, align=PP_ALIGN.CENTER)

    _footer(sl, pg, total)


def _slide_year1(prs, pg, total, ex, result, tariff, proj_df=None,
                 rate_esc_pct=4.0, ppa_cost=None, customer_savings_pct=None):
    """Year 1 cost comparison (left) + savings components chart (right)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _accent_rule(sl)

    bn = result.annual_bill_without_solar
    bw = result.annual_bill_with_solar + (ppa_cost or 0)
    sav = bn - bw
    pct = customer_savings_pct if customer_savings_pct is not None else (sav / bn * 100 if bn else 0)

    _action_title(sl, f"Solar reduces Year 1 electricity cost by {_fd(sav)} ({pct:.1f}%)", exhibit=ex)
    _subtitle(sl, f"Comparison at current {tariff} rates  |  Savings accrue from Day 1")

    # ── LEFT HALF: stacked comparison boxes ──
    lx = ML; lw = Inches(5.0)

    _txt(sl, lx, Inches(1.05), lw, Inches(0.25),
         text="YEAR 1 COST COMPARISON", sz=Pt(11), bold=True, color=DK1)
    _divider_line(sl, Inches(1.33), lw)

    # Without Solar box
    box_w = Inches(2.2); box_h = Inches(1.8)
    bx1 = lx + Inches(0.1); by1 = Inches(1.45)

    _txt(sl, bx1, by1, box_w, Inches(0.20),
         text=f"Without Solar  ({tariff})", sz=Pt(8), color=GRAY50, align=PP_ALIGN.CENTER)
    _rect(sl, bx1, by1 + Inches(0.22), box_w, box_h, fill=DK2)
    _txt(sl, bx1, by1 + Inches(0.22) + box_h * 0.30, box_w, Inches(0.40),
         text=_fd(bn), sz=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, bx1, by1 + Inches(0.22) + box_h * 0.58, box_w, Inches(0.25),
         text="per year", sz=Pt(8), color=GRAY85, align=PP_ALIGN.CENTER)

    # With Solar box
    bx2 = lx + Inches(2.6)
    ratio = bw / bn if bn else 1; sh = box_h * ratio; yo = box_h - sh
    _txt(sl, bx2, by1, box_w, Inches(0.20),
         text="With Solar + Storage", sz=Pt(8), color=GRAY50, align=PP_ALIGN.CENTER)
    _rect(sl, bx2, by1 + Inches(0.22) + yo, box_w, sh, fill=ACCENT1)
    _txt(sl, bx2, by1 + Inches(0.22) + yo + sh * 0.25, box_w, Inches(0.40),
         text=_fd(bw), sz=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, bx2, by1 + Inches(0.22) + yo + sh * 0.55, box_w, Inches(0.25),
         text="per year", sz=Pt(8), color=GRAY85, align=PP_ALIGN.CENTER)

    # Savings callout below boxes
    _txt(sl, lx, Inches(3.65), lw, Inches(0.35),
         text=f"NET SAVINGS:  {_fd(sav)} / YEAR  ({pct:.1f}%)",
         sz=Pt(16), bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

    # Savings breakdown list
    en_sav = result.annual_energy_cost  # with-solar energy
    dm_sav = result.annual_demand_cost  # with-solar demand
    ex_cr = result.annual_export_credit
    # Derive baseline components
    if result.monthly_baseline_details:
        base_en = sum(d.get("energy", 0) for d in result.monthly_baseline_details)
        base_dm = sum(d.get("demand", 0) for d in result.monthly_baseline_details)
    else:
        base_en = bn * 0.55  # fallback estimate
        base_dm = bn * 0.40
    energy_delta = base_en - en_sav
    demand_delta = base_dm - dm_sav

    _txt(sl, lx, Inches(4.10), lw, Inches(0.22),
         text="YEAR 1 SAVINGS BREAKDOWN", sz=Pt(11), bold=True, color=DK1)
    _divider_line(sl, Inches(4.35), lw)

    breakdown = [
        ("Energy charge reduction", _fd(energy_delta), ACCENT1),
        ("Demand charge reduction", _fd(demand_delta), ACCENT3),
        ("Export credits", _fd(ex_cr), ACCENT4),
    ]
    nbc = getattr(result, "annual_nbc_cost", 0) or 0
    if nbc > 0:
        breakdown.append(("NBC charges (NEM-2)", f"({_fd(nbc)})", RGBColor(0xCC, 0x44, 0x44)))

    for i, (lbl, val, clr) in enumerate(breakdown):
        y = Inches(4.47 + i * 0.30)
        _rect(sl, lx + Inches(0.1), y + Inches(0.06), Inches(0.12), Inches(0.12), fill=clr)
        _txt(sl, lx + Inches(0.30), y, Inches(2.8), Inches(0.24),
             text=lbl, sz=Pt(9), color=GRAY50)
        _txt(sl, lx + Inches(3.3), y, Inches(1.5), Inches(0.24),
             text=val, sz=Pt(9), bold=True, color=DK1, align=PP_ALIGN.RIGHT)

    # ── RIGHT HALF: savings components chart ──
    rx = ML + Inches(5.5); rw = Inches(6.3)
    _txt(sl, rx, Inches(1.05), rw, Inches(0.25),
         text="ANNUAL SAVINGS BY COMPONENT", sz=Pt(11), bold=True, color=DK1)
    _divider_line(sl, Inches(1.33), rw, x=rx)

    if proj_df is not None and len(proj_df) > 0 and "Energy ($)" in proj_df.columns:
        # Derive baseline Year 1 components for chart scaling
        if result.monthly_baseline_details:
            yr1_be = sum(d.get("energy", 0) for d in result.monthly_baseline_details)
            yr1_bd = sum(d.get("demand", 0) for d in result.monthly_baseline_details)
            yr1_bf = sum(d.get("fixed", 0) for d in result.monthly_baseline_details)
        else:
            yr1_be = bn * 0.55; yr1_bd = bn * 0.40; yr1_bf = bn * 0.05

        _add_savings_components_chart(sl, rx, Inches(1.40), rw, Inches(3.8),
                                      proj_df, yr1_be, yr1_bd, yr1_bf,
                                      rate_esc_pct)

        # NEM regime annotation if applicable
        nem = getattr(result, "nem_regime", "")
        if nem:
            _txt(sl, rx, Inches(5.30), rw, Inches(0.22),
                 text=f"Current NEM regime: {nem}",
                 sz=Pt(8), italic=True, color=GRAY70)
    else:
        _txt(sl, rx + Inches(0.5), Inches(3.0), Inches(5), Inches(0.30),
             text="Run multi-year projection to populate savings chart",
             sz=Pt(10), italic=True, color=GRAY70, align=PP_ALIGN.CENTER)

    _txt(sl, ML, Inches(5.80), CW, Inches(0.30),
         text="Savings grow each year as utility rates escalate while the solar cost remains contractually fixed.",
         sz=Pt(9), color=GRAY50, align=PP_ALIGN.CENTER)

    _source(sl, "38DN billing simulation at current published tariff rates")
    _footer(sl, pg, total)


def _slide_projections(prs, pg, total, ex, proj_df, rate_esc, ppa=None, esc=None,
                       nem_regime_1=None, nem_regime_2=None, num_years_1=None):
    """Multi-year savings projection table."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _accent_rule(sl)

    df = proj_df.copy()
    if "Year" not in df.columns:
        df.insert(0, "Year", range(1, len(df)+1))
    if "Cumulative Savings ($)" not in df.columns and "Annual Savings ($)" in df.columns:
        df["Cumulative Savings ($)"] = df["Annual Savings ($)"].cumsum()

    cum = df["Cumulative Savings ($)"].iloc[-1] if len(df) else 0
    max_yr = int(df["Year"].max()) if len(df) else 25

    # Determine NEM type per year
    def _nem_label(yr):
        regime = nem_regime_1 or "NEM-3 / NVBT"
        if nem_regime_2 and num_years_1 and yr > num_years_1:
            regime = nem_regime_2
        if "NEM-1" in regime or regime == "NEM-1":
            return "1"
        if "NEM-2" in regime or regime == "NEM-2":
            return "2"
        return "3"

    _action_title(sl, f"Cumulative savings reach {_fd(cum)} over {max_yr} years", exhibit=ex)
    _subtitle(sl, f"Illustrative projection at {rate_esc:.0f}%/yr utility escalation"
              + (f"  |  PPA at ${ppa:.3f}/kWh, {esc:.1f}%/yr escalator" if ppa and esc else ""))
    _takeaway(sl,
        f"At a conservative {rate_esc:.0f}%/yr utility escalation, cumulative savings exceed "
        f"{_fd(cum)}. Actual historical rate growth has been materially higher.")

    hdrs = ["Year", "NEM", "Utility Cost\n(No Solar)", "Total Cost\n(w/ Solar)",
            "Annual\nSavings", "Savings\n%", "Cumulative\nSavings"]
    rows = []
    for _, r in df.iterrows():
        yr = int(r.get("Year", 0))
        ns = r.get("Bill w/o Solar ($)", 0)
        ws = r.get("Bill w/ Solar ($)", 0)
        asv = r.get("Annual Savings ($)", ns - ws)
        p = asv / ns * 100 if ns else 0
        c = r.get("Cumulative Savings ($)", 0)
        rows.append([str(yr), _nem_label(yr), _fd(ns), _fd(ws), _fd(asv), f"{p:.1f}%", _fd(c)])

    cws = [Inches(0.5), Inches(0.5), Inches(1.7), Inches(1.7), Inches(1.5), Inches(0.8), Inches(1.7)]
    _table(sl, ML, Inches(1.80), Inches(8.4), cws, hdrs, rows, bold_last=True, sz=Pt(7))

    note = f"Assumes {rate_esc:.0f}%/yr linear utility rate escalation."
    if ppa and esc:
        note += f" PPA rate of ${ppa:.3f}/kWh at {esc:.1f}%/yr is contractually fixed."
    _source(sl, f"38DN projection model  |  {note}")
    _footer(sl, pg, total)


def _slide_energy(prs, pg, total, ex, result, sys_kw):
    """Monthly energy and billing detail."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _accent_rule(sl)

    offset = result.annual_solar_kwh / result.annual_load_kwh * 100 if result.annual_load_kwh else 0
    yld = result.annual_solar_kwh / sys_kw if sys_kw else 0
    sc = ((result.annual_solar_kwh - result.annual_export_kwh) /
          result.annual_solar_kwh * 100 if result.annual_solar_kwh else 0)

    _action_title(sl, f"System produces {_fk(result.annual_solar_kwh)}/yr, offsetting {offset:.0f}% of site load", exhibit=ex)
    _subtitle(sl, f"Yield: {yld:,.0f} kWh/kW  |  Self-consumption: {sc:.0f}%  |  Export: {_fk(result.annual_export_kwh)}")

    tiles = [
        {"value": _fk(result.annual_solar_kwh), "label": "Annual Production", "accent": ACCENT1},
        {"value": _fk(result.annual_load_kwh), "label": "Annual Load", "accent": ACCENT3},
        {"value": f"{offset:.0f}%", "label": "Solar Offset", "accent": ACCENT1, "val_color": ACCENT1},
        {"value": f"{sc:.0f}%", "label": "Self-Consumption", "accent": ACCENT3},
    ]
    _kpi_row(sl, tiles, Inches(1.15))

    _txt(sl, ML, Inches(2.15), CW, Inches(0.20),
         text="MONTHLY ENERGY & BILLING SUMMARY  (YEAR 1)", sz=Pt(11), bold=True, color=DK1)
    _divider_line(sl, Inches(2.38))

    months = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
    ms = result.monthly_summary
    hdrs = ["Month", "Load\n(MWh)", "Solar\n(MWh)", "Import\n(MWh)", "Export\n(MWh)",
            "Peak\nkW", "Energy\n($)", "Demand\n($)", "Export Cr.\n($)", "Net Bill\n($)"]
    rows = []
    for i in range(min(12, len(ms))):
        r = ms.iloc[i]
        ec = r.get('export_credit', 0)
        rows.append([
            months[i],
            f"{r.get('load_kwh',0)/1e3:,.1f}", f"{r.get('solar_kwh',0)/1e3:,.1f}",
            f"{r.get('import_kwh',0)/1e3:,.1f}", f"{r.get('export_kwh',0)/1e3:,.1f}",
            f"{r.get('peak_demand_kw',0):,.0f}",
            f"${r.get('energy_cost',0):,.0f}", f"${r.get('total_demand_charge',0):,.0f}",
            f"(${ec:,.0f})" if ec > 0 else "$0",
            f"${r.get('net_bill',0):,.0f}",
        ])

    ect = float(ms.get("export_credit", pd.Series([0])).sum())
    rows.append([
        "TOTAL",
        f"{result.annual_load_kwh/1e3:,.1f}", f"{result.annual_solar_kwh/1e3:,.1f}",
        f"{result.annual_import_kwh/1e3:,.1f}", f"{result.annual_export_kwh/1e3:,.1f}",
        f"{ms['peak_demand_kw'].max():,.0f}" if 'peak_demand_kw' in ms.columns else "\u2014",
        f"${result.annual_energy_cost:,.0f}", f"${result.annual_demand_cost:,.0f}",
        f"(${ect:,.0f})" if ect > 0 else "$0",
        f"${result.annual_bill_with_solar:,.0f}",
    ])

    cws = [Inches(0.65)] + [Inches(1.05)] * 4 + [Inches(0.75)] + [Inches(1.05)] * 4
    _table(sl, ML, Inches(2.45), Inches(10.0), cws, hdrs, rows, bold_last=True)

    _source(sl, "38DN 8,760-hour billing simulation  |  All values Year 1")
    _footer(sl, pg, total)


def _slide_production_load(prs, pg, total, ex, result, sys_kw):
    """Monthly production vs load comparison with grouped bar chart."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _accent_rule(sl)

    offset = result.annual_solar_kwh / result.annual_load_kwh * 100 if result.annual_load_kwh else 0
    sc = ((result.annual_solar_kwh - result.annual_export_kwh) /
          result.annual_solar_kwh * 100 if result.annual_solar_kwh else 0)

    _action_title(sl,
        f"Solar offsets {offset:.0f}% of site load ({sc:.0f}% self-consumed)",
        exhibit=ex)
    _subtitle(sl,
        f"{_fk(result.annual_solar_kwh)} annual production  |  "
        f"{_fk(result.annual_load_kwh)} annual load  |  "
        f"{_fk(result.annual_export_kwh)} exported to grid")

    # Section header
    _txt(sl, ML, Inches(1.05), CW, Inches(0.25),
         text="MONTHLY PRODUCTION VS LOAD", sz=Pt(11), bold=True, color=DK1)
    _divider_line(sl, Inches(1.33))

    # Native grouped column chart (Excel-linked, editable in PowerPoint)
    _add_prod_load_chart(sl, ML, Inches(1.45), Inches(10.5), Inches(3.4),
                         result.monthly_summary)

    # KPI row below chart (chart ends at ~4.85")
    yld = result.annual_solar_kwh / sys_kw if sys_kw else 0
    tiles = [
        {"value": f"{offset:.0f}%", "label": "Solar Offset", "accent": ACCENT1, "val_color": ACCENT1},
        {"value": _fk(result.annual_solar_kwh), "label": "Annual Production", "accent": ACCENT1},
        {"value": f"{sc:.0f}%", "label": "Self-Consumption", "accent": ACCENT3},
        {"value": f"{yld:,.0f}", "label": "Specific Yield (kWh/kW)", "accent": DK1},
    ]
    _kpi_row(sl, tiles, Inches(5.10))

    _source(sl, "38DN 8,760-hour simulation; monthly aggregation of hourly production and load profiles")
    _footer(sl, pg, total)


def _slide_rate_hedge(prs, pg, total, ex, utility, ppa, ppa_esc, term,
                      baseline_bill, ppa_bill_yr1, solar_kwh, ppa_rate):
    """Rate hedge analysis with multi-scenario chart."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _accent_rule(sl)
    _action_title(sl,
        "Fixed-rate PPA hedges against rising utility costs",
        exhibit=ex)
    _subtitle(sl, f"{utility}  |  Rate certainty as a strategic advantage")

    _takeaway(sl,
        f"{'The PPA' if ppa else 'Solar'} provides a contractually fixed energy cost for {term} years. "
        f"Day-1 savings are locked in; future utility increases widen the margin.")

    # Native line chart (Excel-linked, editable in PowerPoint)
    _add_hedge_chart(sl, ML, Inches(1.70), Inches(7.5), Inches(4.0),
                     baseline_bill=baseline_bill,
                     ppa_bill_yr1=ppa_bill_yr1,
                     ppa_esc_pct=ppa_esc or 0,
                     term=term,
                     solar_kwh=solar_kwh,
                     ppa_rate_val=ppa_rate)

    # Right side: key points (aligned with chart top)
    rx = ML + Inches(7.8)
    rw = Inches(4.0)
    ry = Inches(1.70)
    _txt(sl, rx, ry, rw, Inches(0.25),
         text="SCENARIO ASSUMPTIONS", sz=Pt(11), bold=True, color=DK1)

    points = [
        f"Middle case: 4%/yr utility escalation",
        f"Upper case: 7%/yr utility escalation",
        f"Stress case: 10%/yr utility escalation",
    ]
    if ppa:
        points.append(f"PPA: ${ppa:.3f}/kWh, {ppa_esc:.1f}%/yr escalator")
    points += [
        "Grid residual escalated at 4%/yr",
        f"Term: {term} years",
    ]
    _bullets(sl, rx + Inches(0.05), ry + Inches(0.40), rw - Inches(0.1),
             points, sz=Pt(10), color=DK1, spacing=Pt(5))

    _txt(sl, rx, Inches(4.20), rw, Inches(0.25),
         text="KEY CONSIDERATIONS", sz=Pt(11), bold=True, color=DK1)

    considerations = [
        "Wildfire liability, grid hardening, and transmission costs drive rates upward",
        "Regulatory cost recovery provides limited customer rate relief",
        (f"PPA rate is contractually fixed for the full {term}-year term"
         if ppa else
         f"Solar cost is fixed for the full {term}-year system life"),
        "Day-1 economics are positive at current rates",
    ]
    _bullets(sl, rx + Inches(0.05), Inches(4.58), rw - Inches(0.1),
             considerations, sz=Pt(10), color=DK1, spacing=Pt(5))

    _source(sl, f"38DN projection model; {utility} published tariff rates; illustrative scenarios")
    _footer(sl, pg, total)


def _slide_process(prs, pg, total, sys_kw, batt_kwh, batt_kw,
                   ppa, esc, utility, tariff, term):
    """Process overview — numbered steps."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _accent_rule(sl)
    _action_title(sl, "Process Overview")
    _subtitle(sl, "From proposal to commercial operation")

    steps = [
        ("System Installation",
         f"{sys_kw:,.0f} kW solar"
         + (f" with {batt_kwh:,.0f} kWh battery" if batt_kwh > 0 else "")
         + " installed on-site at zero cost to you."),
        ("Tariff Optimization",
         f"Transition to the optimal {utility} rate schedule that maximizes value from BTM generation."),
        ("Solar Feeds Operations Directly",
         "Production consumed on-site first, reducing grid imports during daylight hours."
         + (" Excess charges the battery." if batt_kwh > 0 else "")),
    ]
    if batt_kwh > 0:
        steps.append(("Battery Dispatches During Peak Hours",
            f"{batt_kw:,.0f} kW battery discharges during expensive on-peak periods, shaving demand charges."))
    if ppa is not None:
        steps.append(("Predictable, Fixed-Rate Energy",
            f"PPA at ${ppa:.3f}/kWh"
            + (f", {esc:.1f}%/yr escalator" if esc else "")
            + f" for {term} years. No rate volatility on solar kWh."))
    else:
        steps.append(("Immediate Day-1 Savings",
            "Total electricity cost drops from Day 1; gap widens as utility rates rise."))

    # Two columns
    col1 = steps[:3]; col2 = steps[3:]
    for i, (t, d) in enumerate(col1):
        _step(sl, ML + Inches(0.15), Inches(1.30 + i * 1.10), i+1, t, d)
    for i, (t, d) in enumerate(col2):
        _step(sl, ML + Inches(6.3), Inches(1.30 + i * 1.10), len(col1)+i+1, t, d)

    # Ownership bar
    _rect(sl, ML, Inches(5.7), CW, Inches(0.42), fill=DK1)
    _txt(sl, ML + Inches(0.15), Inches(5.74), CW - Inches(0.3), Inches(0.34),
         text=f"System owned, operated, and maintained by 38 Degrees North for the full {term}-year term. "
              f"Zero operational risk or cost to host.",
         sz=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _footer(sl, pg, total)


def _slide_next_steps(prs, pg, total, ppa, esc, term):
    """Next Steps — navy background, green circles."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, Emu(0), Emu(0), SW, SH, fill=DK1)
    _rect(sl, Emu(0), Emu(0), SW, Pt(4), fill=ACCENT1)

    _txt(sl, ML, Inches(0.30), CW, Inches(0.60), text="Next Steps",
         sz=Pt(32), bold=True, color=WHITE)

    steps = [
        ("Execute Letter of Intent",
         "Non-binding LOI confirms mutual interest and key commercial terms."),
        ("Site Assessment & Engineering",
         "Engineering conducts site visit; confirms layout, interconnection, and structural feasibility."),
        ("Execute PPA",
         f"Finalize {term}-year PPA"
         + (f" at ${ppa:.3f}/kWh" if ppa else "")
         + (f" with {esc:.1f}% annual escalator." if esc else ".")),
        ("Permitting & Construction",
         "38DN manages all permits, utility coordination, and construction."),
        ("Commercial Operation",
         "System goes live. Savings begin immediately."),
    ]
    for i, (t, d) in enumerate(steps):
        _step(sl, ML + Inches(0.15), Inches(1.30 + i * 0.98), i+1, t, d,
              circ_c=ACCENT1, dark=True)

    # CTA
    _multi(sl, ML, Inches(6.40), CW, Inches(0.35),
           [{"t": "Lock in rate certainty.  ", "sz": Pt(11), "b": True, "c": ACCENT3},
            {"t": "Zero upfront investment.", "sz": Pt(11), "b": True, "c": ACCENT1}],
           align=PP_ALIGN.CENTER)

    # Wordmark bottom-right
    wp = _logo_path(LOGO_WORD)
    if wp:
        sl.shapes.add_picture(wp, Inches(10.0), Inches(6.85), Inches(2.5))


# ═══════════════════════════════════════════════════════════════════════════
# PUBLIC API
# ═══════════════════════════════════════════════════════════════════════════

def generate_proposal_pptx(
    *,
    customer_name: str,
    address: str = "",
    utility_account: str = "",
    utility_name: str = "",
    tariff_name: str = "",
    new_tariff_name: str | None = None,
    date_str: str = "",
    system_size_kw: float = 0,
    dc_ac_ratio: float = 1.0,
    battery_kwh: float = 0,
    battery_kw: float = 0,
    ppa_rate: float | None = None,
    ppa_escalator_pct: float | None = None,
    term_years: int = 25,
    rate_escalator_pct: float = 4.0,
    result: BillingResult | None = None,
    annual_proj_df: pd.DataFrame | None = None,
    nem_regime_1: str | None = None,
    nem_regime_2: str | None = None,
    num_years_1: int | None = None,
    customer_savings_pct: float | None = None,
    customer_savings_pct_2: float | None = None,
) -> bytes:
    """Generate an institutional-quality branded customer proposal PPTX."""
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH

    has_r = result is not None
    has_p = annual_proj_df is not None and len(annual_proj_df) > 0

    total = 3  # cover + process + next steps
    if has_r: total += 5  # exec summary, current cost, year1, energy detail, production/load
    total += 1  # system design
    if has_p: total += 1  # projections
    total += 1  # rate hedge

    pg = 0; ex = 0

    # Cover
    _slide_cover(prs, customer_name, address, utility_account, date_str, total)

    # Exec Summary
    if has_r:
        pg += 1
        _slide_exec_summary(prs, pg, total, customer_name, result, tariff_name,
                            utility_name, system_size_kw, battery_kwh,
                            ppa_rate, ppa_escalator_pct, term_years,
                            customer_savings_pct=customer_savings_pct,
                            customer_savings_pct_2=customer_savings_pct_2,
                            nem_regime_1=nem_regime_1,
                            nem_regime_2=nem_regime_2)

    # Current Cost
    if has_r:
        pg += 1; ex += 1
        _slide_current_cost(prs, pg, total, ex, customer_name, result,
                            tariff_name, utility_name)

    # System Design
    pg += 1; ex += 1
    _slide_system(prs, pg, total, ex, system_size_kw, dc_ac_ratio,
                  battery_kwh, battery_kw, ppa_rate, ppa_escalator_pct,
                  term_years, tariff_name, new_tariff_name)

    # Year 1 Comparison
    if has_r:
        pg += 1; ex += 1
        _slide_year1(prs, pg, total, ex, result, tariff_name,
                     proj_df=annual_proj_df, rate_esc_pct=rate_escalator_pct,
                     customer_savings_pct=customer_savings_pct)

    # Projections
    if has_p:
        pg += 1; ex += 1
        _slide_projections(prs, pg, total, ex, annual_proj_df,
                           rate_escalator_pct, ppa_rate, ppa_escalator_pct,
                           nem_regime_1=nem_regime_1, nem_regime_2=nem_regime_2,
                           num_years_1=num_years_1)

    # Production vs Load
    if has_r:
        pg += 1; ex += 1
        _slide_production_load(prs, pg, total, ex, result, system_size_kw)

    # Energy Detail
    if has_r:
        pg += 1; ex += 1
        _slide_energy(prs, pg, total, ex, result, system_size_kw)

    # Process
    pg += 1
    _slide_process(prs, pg, total, system_size_kw, battery_kwh, battery_kw,
                   ppa_rate, ppa_escalator_pct, utility_name, tariff_name, term_years)

    # Rate Hedge
    pg += 1; ex += 1
    _slide_rate_hedge(prs, pg, total, ex, utility_name,
                      ppa_rate, ppa_escalator_pct, term_years,
                      baseline_bill=result.annual_bill_without_solar if has_r else 0,
                      ppa_bill_yr1=result.annual_bill_with_solar if has_r else 0,
                      solar_kwh=result.annual_solar_kwh if has_r else 0,
                      ppa_rate=ppa_rate)

    # Next Steps
    pg += 1
    _slide_next_steps(prs, pg, total, ppa_rate, ppa_escalator_pct, term_years)

    buf = BytesIO(); prs.save(buf); return buf.getvalue()
